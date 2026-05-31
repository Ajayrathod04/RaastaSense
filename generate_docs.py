import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

import docx
from docx.shared import Inches as DocxInches, Pt as DocxPt, RGBColor as DocxRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml import OxmlElement, parse_xml
from docx.oxml.ns import nsdecls, qn

def apply_background(slide, color_rgb):
    """Fills the slide background with a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color_rgb

def create_presentation():
    print("Generating PowerPoint Presentation...")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 🎨 Brand Palette
    BG_DARK = RGBColor(15, 23, 42)       # #0F172A
    SKY_BLUE = RGBColor(56, 189, 248)    # #38BDF8
    CRIMSON_RED = RGBColor(244, 63, 94)  # #F43F5E
    AMBER_GOLD = RGBColor(245, 158, 11)  # #F59E0B
    WHITE = RGBColor(255, 255, 255)
    MUTED_SLATE = RGBColor(148, 163, 184) # #94A3B8
    CARD_BG = RGBColor(30, 41, 59)       # #1E293B

    blank_layout = prs.slide_layouts[6] # Blank slide layout

    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    apply_background(slide1, BG_DARK)
    
    # Left accent pillar shape
    shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(0.15), Inches(6.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SKY_BLUE
    shape.line.fill.background()

    # Main text box
    tb = slide1.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(11.5), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p_title = tf.paragraphs[0]
    p_title.text = "RaastaSense"
    p_title.font.name = "Segoe UI"
    p_title.font.size = Pt(64)
    p_title.font.bold = True
    p_title.font.color.rgb = WHITE
    p_title.space_after = Pt(10)
    
    p_tag = tf.add_paragraph()
    p_tag.text = "Transforming Road Safety Through Intelligence"
    p_tag.font.name = "Segoe UI"
    p_tag.font.size = Pt(22)
    p_tag.font.bold = True
    p_tag.font.color.rgb = SKY_BLUE
    p_tag.space_after = Pt(50)
    
    p_team_name = tf.add_paragraph()
    p_team_name.text = "Team: NXT-waVE (IIT Madras Road Safety Hackathon)"
    p_team_name.font.name = "Segoe UI"
    p_team_name.font.size = Pt(14)
    p_team_name.font.bold = True
    p_team_name.font.color.rgb = AMBER_GOLD
    
    p_members = tf.add_paragraph()
    p_members.text = "AJAY RATHOD  |  ANIMESH BHANARKAR"
    p_members.font.name = "Segoe UI"
    p_members.font.size = Pt(16)
    p_members.font.bold = True
    p_members.font.color.rgb = WHITE
    
    # ==========================================
    # SLIDE 2: Problem Statement
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    apply_background(slide2, BG_DARK)
    
    # Title
    tb_title = slide2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
    tf_title = tb_title.text_frame
    p = tf_title.paragraphs[0]
    p.text = "THE NATIONAL ROAD SAFETY CRISIS"
    p.font.name = "Segoe UI"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = CRIMSON_RED
    
    # Left Card: Statistics
    card_l = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(4.8))
    card_l.fill.solid()
    card_l.fill.fore_color.rgb = CARD_BG
    card_l.line.color.rgb = CRIMSON_RED
    card_l.line.width = Pt(1.5)
    
    tf_l = card_l.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = Inches(0.3)
    tf_l.margin_right = Inches(0.3)
    tf_l.margin_top = Inches(0.3)
    
    p_l1 = tf_l.paragraphs[0]
    p_l1.text = "🚨 Critical National Statistics"
    p_l1.font.name = "Segoe UI"
    p_l1.font.size = Pt(20)
    p_l1.font.bold = True
    p_l1.font.color.rgb = WHITE
    p_l1.space_after = Pt(20)
    
    p_l2 = tf_l.add_paragraph()
    p_l2.text = "• 1.5 Lakhs+ Annual Fatalities\n  Indian corridors average over 400 deaths per day.\n\n• ₹5.8 Lakh Crores Economic Loss\n  Accidents deplete ~3% of India's annual GDP.\n\n• 50%+ Preventable Deaths\n  Delayed Golden Hour response & static navigation."
    p_l2.font.name = "Segoe UI"
    p_l2.font.size = Pt(14)
    p_l2.font.color.rgb = MUTED_SLATE
    
    # Right Card: Pain Points
    card_r = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.6), Inches(4.8))
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = CARD_BG
    card_r.line.color.rgb = SKY_BLUE
    card_r.line.width = Pt(1.5)
    
    tf_r = card_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = Inches(0.3)
    tf_r.margin_right = Inches(0.3)
    tf_r.margin_top = Inches(0.3)
    
    p_r1 = tf_r.paragraphs[0]
    p_r1.text = "⚠️ Core Infrastructure Gaps"
    p_r1.font.name = "Segoe UI"
    p_r1.font.size = Pt(20)
    p_r1.font.bold = True
    p_r1.font.color.rgb = WHITE
    p_r1.space_after = Pt(20)
    
    p_r2 = tf_r.add_paragraph()
    p_r2.text = "• Responders blind to real-time satellite coordinates & golden hour locks.\n\n• Civic reports lacking public audit trails or visual pothole classification.\n\n• Driving violations static across municipal jurisdictions with no localized counselor."
    p_r2.font.name = "Segoe UI"
    p_r2.font.size = Pt(14)
    p_r2.font.color.rgb = MUTED_SLATE

    # ==========================================
    # SLIDE 3: Solution
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    apply_background(slide3, BG_DARK)
    
    # Title
    tb_title = slide3.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
    tf_title = tb_title.text_frame
    p = tf_title.paragraphs[0]
    p.text = "THE RAASTASENSE SOLUTION"
    p.font.name = "Segoe UI"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = SKY_BLUE
    
    # Pillar 1
    p1 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.5))
    p1.fill.solid()
    p1.fill.fore_color.rgb = CARD_BG
    p1.line.color.rgb = SKY_BLUE
    tf1 = p1.text_frame
    tf1.word_wrap = True
    tf1.margin_top = Inches(0.3)
    p_t1 = tf1.paragraphs[0]
    p_t1.text = "🗺️ GIS Live Map HUD"
    p_t1.font.name = "Segoe UI"
    p_t1.font.size = Pt(18)
    p_t1.font.bold = True
    p_t1.font.color.rgb = WHITE
    p_t1.space_after = Pt(14)
    p_d1 = tf1.add_paragraph()
    p_d1.text = "Cyberslate interactive mapping featuring active incident reporting overlays, speedometers, and live coordinate locks."
    p_d1.font.name = "Segoe UI"
    p_d1.font.size = Pt(13)
    p_d1.font.color.rgb = MUTED_SLATE
    
    # Pillar 2
    p2 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.5))
    p2.fill.solid()
    p2.fill.fore_color.rgb = CARD_BG
    p2.line.color.rgb = AMBER_GOLD
    tf2 = p2.text_frame
    tf2.word_wrap = True
    tf2.margin_top = Inches(0.3)
    p_t2 = tf2.paragraphs[0]
    p_t2.text = "🛡️ AI Legal Advisor"
    p_t2.font.name = "Segoe UI"
    p_t2.font.size = Pt(18)
    p_t2.font.bold = True
    p_t2.font.color.rgb = WHITE
    p_t2.space_after = Pt(14)
    p_d2 = tf2.add_paragraph()
    p_d2.text = "9-language custom localizer featuring automatic fine surcharges, legal index, and character dialogue advisory boards."
    p_d2.font.name = "Segoe UI"
    p_d2.font.size = Pt(13)
    p_d2.font.color.rgb = MUTED_SLATE

    # Pillar 3
    p3 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.5))
    p3.fill.solid()
    p3.fill.fore_color.rgb = CARD_BG
    p3.line.color.rgb = CRIMSON_RED
    tf3 = p3.text_frame
    tf3.word_wrap = True
    tf3.margin_top = Inches(0.3)
    p_t3 = tf3.paragraphs[0]
    p_t3.text = "🆘 SOS Rescue Hub"
    p_t3.font.name = "Segoe UI"
    p_t3.font.size = Pt(18)
    p_t3.font.bold = True
    p_t3.font.color.rgb = WHITE
    p_t3.space_after = Pt(14)
    p_d3 = tf3.add_paragraph()
    p_d3.text = "Decentralized emergency countdown beacon, golden-hour responder coordination, and pre-routed medical pathways."
    p_d3.font.name = "Segoe UI"
    p_d3.font.size = Pt(13)
    p_d3.font.color.rgb = MUTED_SLATE

    # ==========================================
    # SLIDE 4: Key Features (Core Engines)
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    apply_background(slide4, BG_DARK)
    
    # Title
    tb_title = slide4.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
    tf_title = tb_title.text_frame
    p = tf_title.paragraphs[0]
    p.text = "INTELLIGENT CORE ENGINE MODULES"
    p.font.name = "Segoe UI"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = AMBER_GOLD
    
    # 2x2 Widescreen Grid
    # Card 1: Safe Route Engine
    c1 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(2.2))
    c1.fill.solid()
    c1.fill.fore_color.rgb = CARD_BG
    c1.line.color.rgb = SKY_BLUE
    t_c1 = c1.text_frame
    t_c1.word_wrap = True
    t_c1.margin_top = Inches(0.15)
    t_c1.paragraphs[0].text = "🚦 Safe Route Engine (Phase 8 Highlight)"
    t_c1.paragraphs[0].font.name = "Segoe UI"
    t_c1.paragraphs[0].font.size = Pt(16)
    t_c1.paragraphs[0].font.bold = True
    t_c1.paragraphs[0].font.color.rgb = WHITE
    t_c1.paragraphs[0].space_after = Pt(8)
    d_c1 = t_c1.add_paragraph()
    d_c1.text = "Calculates Fastest, Safest, and Least Incident corridors. Displays Safety Indexes, weather/traffic conditions, and dynamic recommended overlays."
    d_c1.font.name = "Segoe UI"
    d_c1.font.size = Pt(12)
    d_c1.font.color.rgb = MUTED_SLATE
    
    # Card 2: DriveLegal Surcharge Sorter
    c2 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.6), Inches(2.2))
    c2.fill.solid()
    c2.fill.fore_color.rgb = CARD_BG
    c2.line.color.rgb = SKY_BLUE
    t_c2 = c2.text_frame
    t_c2.word_wrap = True
    t_c2.margin_top = Inches(0.15)
    t_c2.paragraphs[0].text = "⚖️ DriveLegal Regional Surcharges"
    t_c2.paragraphs[0].font.name = "Segoe UI"
    t_c2.paragraphs[0].font.size = Pt(16)
    t_c2.paragraphs[0].font.bold = True
    t_c2.paragraphs[0].font.color.rgb = WHITE
    t_c2.paragraphs[0].space_after = Pt(8)
    d_c2 = t_c2.add_paragraph()
    d_c2.text = "Filters 6 main road safety codes. Instantly computes fine adjustments based on Indian State surcharges (TN, MH, KA) and delivers voice-synthesized guidelines."
    d_c2.font.name = "Segoe UI"
    d_c2.font.size = Pt(12)
    d_c2.font.color.rgb = MUTED_SLATE

    # Card 3: RoadWatch AI Pothole Diagnosis
    c3 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.2), Inches(5.6), Inches(2.2))
    c3.fill.solid()
    c3.fill.fore_color.rgb = CARD_BG
    c3.line.color.rgb = SKY_BLUE
    t_c3 = c3.text_frame
    t_c3.word_wrap = True
    t_c3.margin_top = Inches(0.15)
    t_c3.paragraphs[0].text = "📸 RoadWatch AI Incident Diagnostics"
    t_c3.paragraphs[0].font.name = "Segoe UI"
    t_c3.paragraphs[0].font.size = Pt(16)
    t_c3.paragraphs[0].font.bold = True
    t_c3.paragraphs[0].font.color.rgb = WHITE
    t_c3.paragraphs[0].space_after = Pt(8)
    d_c3 = t_c3.add_paragraph()
    d_c3.text = "Citizen hazard submission portals. Features automatic AI-powered image diagnosis with calculated confidence matrices, severity metrics, and direct action protocols."
    d_c3.font.name = "Segoe UI"
    d_c3.font.size = Pt(12)
    d_c3.font.color.rgb = MUTED_SLATE

    # Card 4: RoadSOS Emergency Dispatch
    c4 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(4.2), Inches(5.6), Inches(2.2))
    c4.fill.solid()
    c4.fill.fore_color.rgb = CARD_BG
    c4.line.color.rgb = SKY_BLUE
    t_c4 = c4.text_frame
    t_c4.word_wrap = True
    t_c4.margin_top = Inches(0.15)
    t_c4.paragraphs[0].text = "🆘 RoadSOS Emergency Controller"
    t_c4.paragraphs[0].font.name = "Segoe UI"
    t_c4.paragraphs[0].font.size = Pt(16)
    t_c4.paragraphs[0].font.bold = True
    t_c4.paragraphs[0].font.color.rgb = WHITE
    t_c4.paragraphs[0].space_after = Pt(8)
    d_c4 = t_c4.add_paragraph()
    d_c4.text = "GPS telemetry tracker locks precise coordinates. Displays immediate emergency dials, trauma center proximities, and custom full-screen alarm dispatch overlay panels."
    d_c4.font.name = "Segoe UI"
    d_c4.font.size = Pt(12)
    d_c4.font.color.rgb = MUTED_SLATE

    # ==========================================
    # SLIDE 5: System Architecture & HUD
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    apply_background(slide5, BG_DARK)
    
    # Title
    tb_title = slide5.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
    tf_title = tb_title.text_frame
    p = tf_title.paragraphs[0]
    p.text = "SYSTEM ARCHITECTURE & OBSERVABILITY HUD"
    p.font.name = "Segoe UI"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Left Card: Technology Stack & Ingestion Pipeline
    arch_card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(4.8))
    arch_card.fill.solid()
    arch_card.fill.fore_color.rgb = CARD_BG
    arch_card.line.color.rgb = SKY_BLUE
    arch_card.line.width = Pt(1.5)
    tf_arch = arch_card.text_frame
    tf_arch.word_wrap = True
    tf_arch.margin_top = Inches(0.3)
    p_at = tf_arch.paragraphs[0]
    p_at.text = "⚙️ Full-Stack Pipeline Architecture"
    p_at.font.name = "Segoe UI"
    p_at.font.size = Pt(20)
    p_at.font.bold = True
    p_at.font.color.rgb = WHITE
    p_at.space_after = Pt(18)
    
    p_ad = tf_arch.add_paragraph()
    p_ad.text = "• Client Tier: React 18, Vite, Tailwind CSS, Framer Motion, and custom vector SVGs.\n\n• Map/GIS Layer: OpenStreetMap integrations, polyline vector algorithms, dynamic marker overlays.\n\n• Backend & Database: Node.js, Express, mock-telemetry sensor algorithms.\n\n• AI Core: Google Vision API pothole pixel scanner & safety chatbot fallbacks."
    p_ad.font.name = "Segoe UI"
    p_ad.font.size = Pt(13)
    p_ad.font.color.rgb = MUTED_SLATE

    # Right Card: Observability & Telemetry HUD
    hud_card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.6), Inches(4.8))
    hud_card.fill.solid()
    hud_card.fill.fore_color.rgb = CARD_BG
    hud_card.line.color.rgb = AMBER_GOLD
    hud_card.line.width = Pt(1.5)
    tf_hud = hud_card.text_frame
    tf_hud.word_wrap = True
    tf_hud.margin_top = Inches(0.3)
    p_ht = tf_hud.paragraphs[0]
    p_ht.text = "📊 Live Observability & Telemetry Console"
    p_ht.font.name = "Segoe UI"
    p_ht.font.size = Pt(20)
    p_ht.font.bold = True
    p_ht.font.color.rgb = WHITE
    p_ht.space_after = Pt(18)
    
    p_hd = tf_hud.add_paragraph()
    p_hd.text = "• Live Console Telemetry Stream: Chronological NoSQL Firestore connectivity indicator and telemetry status outputs.\n\n• Custom SVG Safety Charts: Real-time route congestion graph and repairs metrics tracker.\n\n• 9-Language Localization Dropdown: Multi-lingual localizer supporting regional fine surcharges, legal codes, and AI advice swaps with zero-lag."
    p_hd.font.name = "Segoe UI"
    p_hd.font.size = Pt(13)
    p_hd.font.color.rgb = MUTED_SLATE

    # ==========================================
    # SLIDE 6: Societal Impact
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    apply_background(slide6, BG_DARK)
    
    # Title
    tb_title = slide6.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
    tf_title = tb_title.text_frame
    p = tf_title.paragraphs[0]
    p.text = "MEASURED SOCIETAL & CIVIC IMPACT"
    p.font.name = "Segoe UI"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = SKY_BLUE
    
    # 3 Column split
    w_col = Inches(3.6)
    # Col 1: Citizens
    c1_s = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), w_col, Inches(4.5))
    c1_s.fill.solid()
    c1_s.fill.fore_color.rgb = CARD_BG
    c1_s.line.color.rgb = SKY_BLUE
    tf_c1 = c1_s.text_frame
    tf_c1.word_wrap = True
    tf_c1.margin_top = Inches(0.3)
    p_col_1 = tf_c1.paragraphs[0]
    p_col_1.text = "🚲 Citizens"
    p_col_1.font.name = "Segoe UI"
    p_col_1.font.size = Pt(18)
    p_col_1.font.bold = True
    p_col_1.font.color.rgb = WHITE
    p_col_1.space_after = Pt(14)
    p_desc_1 = tf_c1.add_paragraph()
    p_desc_1.text = "Safe routing guidelines keep travelers out of congestion/danger zones. Regional DriveLegal counselor builds critical legal self-awareness."
    p_desc_1.font.name = "Segoe UI"
    p_desc_1.font.size = Pt(13)
    p_desc_1.font.color.rgb = MUTED_SLATE
    
    # Col 2: Authorities
    c2_s = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(1.8), w_col, Inches(4.5))
    c2_s.fill.solid()
    c2_s.fill.fore_color.rgb = CARD_BG
    c2_s.line.color.rgb = AMBER_GOLD
    tf_c2 = c2_s.text_frame
    tf_c2.word_wrap = True
    tf_c2.margin_top = Inches(0.3)
    p_col_2 = tf_c2.paragraphs[0]
    p_col_2.text = "👮 Municipalities"
    p_col_2.font.name = "Segoe UI"
    p_col_2.font.size = Pt(18)
    p_col_2.font.bold = True
    p_col_2.font.color.rgb = WHITE
    p_col_2.space_after = Pt(14)
    p_desc_2 = tf_c2.add_paragraph()
    p_desc_2.text = "RoadWatch community dashboards create transparency in public infrastructure spending, pothole repair trails, and automated violation reporting."
    p_desc_2.font.name = "Segoe UI"
    p_desc_2.font.size = Pt(13)
    p_desc_2.font.color.rgb = MUTED_SLATE

    # Col 3: Responders
    c3_s = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.8), w_col, Inches(4.5))
    c3_s.fill.solid()
    c3_s.fill.fore_color.rgb = CARD_BG
    c3_s.line.color.rgb = CRIMSON_RED
    tf_c3 = c3_s.text_frame
    tf_c3.word_wrap = True
    tf_c3.margin_top = Inches(0.3)
    p_col_3 = tf_c3.paragraphs[0]
    p_col_3.text = "🚑 Responders"
    p_col_3.font.name = "Segoe UI"
    p_col_3.font.size = Pt(18)
    p_col_3.font.bold = True
    p_col_3.font.color.rgb = WHITE
    p_col_3.space_after = Pt(14)
    p_desc_3 = tf_c3.add_paragraph()
    p_desc_3.text = "RoadSOS beacon coordinates lock directly on responder navigation units, slicing dispatch delays and securing the critical Golden Hour rescue frame."
    p_desc_3.font.name = "Segoe UI"
    p_desc_3.font.size = Pt(13)
    p_desc_3.font.color.rgb = MUTED_SLATE

    # ==========================================
    # SLIDE 7: Thank You & Resources
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    apply_background(slide7, BG_DARK)
    
    tb_thank = slide7.shapes.add_textbox(Inches(0.9), Inches(1.2), Inches(11.5), Inches(5.0))
    tf_thank = tb_thank.text_frame
    tf_thank.word_wrap = True
    
    p_thank = tf_thank.paragraphs[0]
    p_thank.text = "THANK YOU!"
    p_thank.alignment = PP_ALIGN.CENTER
    p_thank.font.name = "Segoe UI"
    p_thank.font.size = Pt(64)
    p_thank.font.bold = True
    p_thank.font.color.rgb = SKY_BLUE
    p_thank.space_after = Pt(14)
    
    p_sub = tf_thank.add_paragraph()
    p_sub.text = "Join us in transforming Indian transit corridors, one life at a time."
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.font.name = "Segoe UI"
    p_sub.font.size = Pt(20)
    p_sub.font.bold = True
    p_sub.font.color.rgb = WHITE
    p_sub.space_after = Pt(40)
    
    p_res = tf_thank.add_paragraph()
    p_res.text = "🌐 Interactive Live Demo: https://raastasense.vercel.app\n💻 Open-Source Codebase: https://github.com/Ajayrathod04/RaastaSense"
    p_res.alignment = PP_ALIGN.CENTER
    p_res.font.name = "Consolas"
    p_res.font.size = Pt(16)
    p_res.font.bold = True
    p_res.font.color.rgb = AMBER_GOLD
    p_res.space_after = Pt(20)

    prs.save("RaastaSense_Hackathon_Presentation.pptx")
    print("PowerPoint presentation generated successfully!")

# XML helper elements for Styling docx
def set_cell_background(cell, fill_hex):
    tcPr = cell._tc.get_or_add_tcPr()
    shd = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{fill_hex}"/>')
    tcPr.append(shd)

def set_cell_margins(cell, top=100, bottom=100, left=150, right=150):
    tcPr = cell._tc.get_or_add_tcPr()
    tcMar = parse_xml(f'<w:tcMar {nsdecls("w")}><w:top w:w="{top}" w:type="dxa"/><w:bottom w:w="{bottom}" w:type="dxa"/><w:left w:w="{left}" w:type="dxa"/><w:right w:w="{right}" w:type="dxa"/></w:tcMar>')
    tcPr.append(tcMar)

def create_docx_proposal():
    print("Generating DOCX Hackathon Proposal...")
    doc = docx.Document()
    
    # Page setup
    sections = doc.sections
    for section in sections:
        section.top_margin = DocxInches(1)
        section.bottom_margin = DocxInches(1)
        section.left_margin = DocxInches(1)
        section.right_margin = DocxInches(1)

    # 🎨 Color Schemes
    C_BLUE = DocxRGBColor(56, 189, 248) # #38BDF8
    C_GOLD = DocxRGBColor(245, 158, 11) # #F59E0B
    C_SLATE = DocxRGBColor(15, 23, 42) # #0F172A
    C_GRAY = DocxRGBColor(71, 85, 105) # #475569

    # Setup styles
    style_normal = doc.styles['Normal']
    style_normal.font.name = 'Calibri'
    style_normal.font.size = DocxPt(11)
    style_normal.font.color.rgb = C_SLATE
    style_normal.paragraph_format.line_spacing = 1.15
    style_normal.paragraph_format.space_after = DocxPt(6)

    # Helper functions for titles
    def add_chapter_title(text):
        p = doc.add_paragraph()
        run = p.add_run(text)
        run.font.name = 'Century Gothic'
        run.font.size = DocxPt(20)
        run.font.bold = True
        run.font.color.rgb = C_SLATE
        p.paragraph_format.space_before = DocxPt(18)
        p.paragraph_format.space_after = DocxPt(10)
        return p

    def add_section_heading(text):
        p = doc.add_paragraph()
        run = p.add_run(text)
        run.font.name = 'Century Gothic'
        run.font.size = DocxPt(14)
        run.font.bold = True
        run.font.color.rgb = C_GRAY
        p.paragraph_format.space_before = DocxPt(12)
        p.paragraph_format.space_after = DocxPt(6)
        return p

    # ==========================================
    # CHAPTER 1: COVER PAGE
    # ==========================================
    for _ in range(3): doc.add_paragraph()
    
    p_title = doc.add_paragraph()
    p_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run_t = p_title.add_run("RAASTASENSE V2")
    run_t.font.name = 'Century Gothic'
    run_t.font.size = DocxPt(44)
    run_t.font.bold = True
    run_t.font.color.rgb = C_SLATE
    p_title.paragraph_format.space_after = DocxPt(6)

    p_sub = doc.add_paragraph()
    p_sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run_sub = p_sub.add_run("Transforming Road Safety Through Real-Time Intelligence")
    run_sub.font.name = 'Century Gothic'
    run_sub.font.size = DocxPt(16)
    run_sub.font.bold = True
    run_sub.font.color.rgb = C_BLUE
    p_sub.paragraph_format.space_after = DocxPt(40)

    for _ in range(4): doc.add_paragraph()

    # Meta box
    p_meta = doc.add_paragraph()
    p_meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run_meta = p_meta.add_run(
        "IIT Madras Road Safety Hackathon 2026\n"
        "Team: NXT-waVE\n"
        "Members: AJAY RATHOD & ANIMESH BHANARKAR\n"
        "Live Deployment: https://raastasense.vercel.app\n"
        "Repository: https://github.com/Ajayrathod04/RaastaSense\n"
    )
    run_meta.font.name = 'Calibri'
    run_meta.font.size = DocxPt(12)
    run_meta.font.bold = True
    run_meta.font.color.rgb = C_GRAY
    
    doc.add_page_break()

    # ==========================================
    # CHAPTER 2: EXECUTIVE SUMMARY
    # ==========================================
    add_chapter_title("1. Executive Summary")
    doc.add_paragraph(
        "RaastaSense V2 represents a high-fidelity paradigm shift in modern road safety management. Designed specifically for "
        "the IIT Madras Road Safety Challenge, this ecosystem integrates reactive public governance trackers, predictive GIS safe route navigation, "
        "and direct satellite GPS responder coordination. Over half of road fatalities in developing nations occur during the critical 'Golden Hour' "
        "due to dispatcher routing delays and public blind spots. RaastaSense V2 addresses these challenges by transforming community vigilance "
        "into a highly transparent, actionable dashboard. With full 9-language i18n support, state-specific fine surcharges filter nodes, automated "
        "AI incident pixel categorization models, and deep coordinate locks, the platform bridges the gap between citizens, municipal agencies, "
        "and medical emergency teams."
    )
    doc.add_page_break()

    # ==========================================
    # CHAPTER 3: PROBLEM STATEMENT
    # ==========================================
    add_chapter_title("2. Problem Statement")
    doc.add_paragraph(
        "Indian transit corridors are suffering from severe structural inefficiencies, contributing to significant loss of life and financial strain. "
        "Critical pain points include:"
    )
    doc.add_paragraph("• High Annual Fatalities: Indian corridors suffer from more than 1.5 Lakhs fatalities annually, averaging 400 lives lost every single day.")
    doc.add_paragraph("• Massive Economic Loss: Road traffic accidents drain approximately ₹5.8 Lakh Crores annually, capturing almost 3% of national GDP.")
    doc.add_paragraph("• Golden Hour Responding: Dispatchers struggle with delayed responses because of poor location telemetry, routing bottlenecks, and a lack of real-time coordinate locks.")
    doc.add_paragraph("• Infrastructure Transparency: Hazard spots like severe potholes remain unlogged and unresolved due to a lack of public audit trails or visual pothole classification pipelines.")
    doc.add_paragraph("• Dynamic Surcharges: Legal violations indexes remain static, failing to adapt to dynamic municipal surcharge changes (such as the distinct fines between Tamil Nadu, Maharashtra, or Karnataka).")
    doc.add_page_break()

    # ==========================================
    # CHAPTER 4: PROPOSED SOLUTION
    # ==========================================
    add_chapter_title("3. Proposed Solution")
    doc.add_paragraph(
        "RaastaSense V2 introduces an interactive, fail-safe AI Road Safety Ecosystem. Key features include:"
    )
    doc.add_paragraph("1. Cyberslate Live GIS HUD: Direct interactive Leaflet mapping displaying real-time speedometers, traffic-fluctuation tile layers, coordinate markers, and live telemetry simulators.")
    doc.add_paragraph("2. DriveLegal Indexer: Dynamically filters 6 core legal safety codes (overspeeding, mobile usage, helmets) with state-specific surcharge modifiers and voice synthesizers.")
    doc.add_paragraph("3. RoadWatch AI Audits: Transparent complaint logs paired with an automatic AI incident diagnostics console. Mock images submitted receive an immediate, multi-tier pixel scan (Severity Index, Risk Category, and Suggested Action Protocols).")
    doc.add_paragraph("4. RoadSOS Emergency Beacons: Sat-locked coordinates, trauma center proximity indexing, direct dial numbers, and a customized full-screen Emergency SOS alarm override panel that routes emergency dispatchers instantly.")

    # Screenshot 1
    doc.add_paragraph()
    if os.path.exists("raastasense_hero_ui.png"):
        doc.add_picture("raastasense_hero_ui.png", width=DocxInches(6.0))
        p_cap = doc.add_paragraph()
        p_cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run_cap = p_cap.add_run("Figure 1: RaastaSense V2 Glassmorphic Cyber Intelligence Dashboard")
        run_cap.font.italic = True
        run_cap.font.size = DocxPt(9.5)
        run_cap.font.color.rgb = C_GRAY
    doc.add_page_break()

    # ==========================================
    # CHAPTER 5: SYSTEM ARCHITECTURE
    # ==========================================
    add_chapter_title("4. System Architecture")
    doc.add_paragraph(
        "The RaastaSense platform utilizes a highly scalable, decoupled client-server architecture. The components include:"
    )
    doc.add_paragraph("• Client Tier: React 18, Vite, Framer Motion, custom vector SVGs, and Tailwind CSS.")
    doc.add_paragraph("• Map Layer: OpenStreetMap integrations, polyline coordinates generator, and dynamic marker overlays.")
    doc.add_paragraph("• Backend & API: Express server, telemetry polling, and NoSQL Firestore mock logs.")
    doc.add_paragraph("• AI Diagnostics: Pothole pixel classification model using Google Generative AI (Gemini) fallback channels.")
    doc.add_page_break()

    # ==========================================
    # CHAPTER 6: CORE ENGINE MODULES
    # ==========================================
    add_chapter_title("5. Core Engine Modules")
    
    add_section_heading("5.1 Safe Route Engine (Phase 8 Highlight)")
    doc.add_paragraph(
        "The dynamic navigation module generates three distinct routing corridors (Fastest, Safest, and Least Incident path). "
        "It visualizes safety percentages, traffic congestion indexes, accident probability ratios, and rain/weather impact states, "
        "applying distinct highlights to emphasize the absolute safest route for two-wheeler or four-wheeler safety."
    )
    # Screenshot 2
    if os.path.exists("raastasense_route_engine.png"):
        doc.add_picture("raastasense_route_engine.png", width=DocxInches(6.0))
        p_cap = doc.add_paragraph()
        p_cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run_cap = p_cap.add_run("Figure 2: Safe Route Corridor Navigation HUD & Deep Safety Metrics")
        run_cap.font.italic = True
        run_cap.font.size = DocxPt(9.5)
        run_cap.font.color.rgb = C_GRAY

    add_section_heading("5.2 DriveLegal Index")
    doc.add_paragraph(
        "Fosters citizen awareness by displaying detailed explanations of violations. Standard fines scale dynamically based on selected "
        "Indian state rules (e.g. 20% fine surcharge in Tamil Nadu, 30% in Maharashtra) to simulate local transport department parameters."
    )

    add_section_heading("5.3 RoadWatch Incident Diagnostics")
    doc.add_paragraph(
        "Allows citizens to upload image reports. An automated, simulated AI pixel diagnostics console executes a multi-layer evaluation "
        "indexing severity out of 10, confidence rates, hazard categorization, and direct action guidelines (e.g. Immediate Asphalt Patching)."
    )

    add_section_heading("5.4 RoadSOS Emergency Controller")
    doc.add_paragraph(
        "Emergency coordinate resolve logs direct sat-locks. Shows responder proximities (Hospitals & Police centers) with one-click dial anchors, "
        "and triggers a beautiful full-screen overlay panel showing evacuation escape routes and immediate dispatcher twilio broadcasts."
    )
    # Screenshot 3
    if os.path.exists("raastasense_observability_settings.png"):
        doc.add_picture("raastasense_observability_settings.png", width=DocxInches(6.0))
        p_cap = doc.add_paragraph()
        p_cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run_cap = p_cap.add_run("Figure 3: System HUD, Multilingual Localization and Live Cloud Telemetry Console")
        run_cap.font.italic = True
        run_cap.font.size = DocxPt(9.5)
        run_cap.font.color.rgb = C_GRAY
    doc.add_page_break()

    # ==========================================
    # CHAPTER 7: TECHNOLOGY STACK & DEPENDENCIES
    # ==========================================
    add_chapter_title("6. Technology Stack & Software Packages")
    doc.add_paragraph(
        "Below is a tabulated summary of the technologies and core packages deployed within RaastaSense V2:"
    )

    table = doc.add_table(rows=6, cols=3)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.autofit = False
    
    headers = ["Component", "Technologies / Libraries", "Purpose"]
    col_widths = [DocxInches(1.8), DocxInches(2.5), DocxInches(2.2)]
    
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.paragraphs[0].runs[0].font.bold = True
        cell.paragraphs[0].runs[0].font.color.rgb = DocxRGBColor(255, 255, 255)
        set_cell_background(cell, "0F172A")
        set_cell_margins(cell)

    tech_data = [
        ["Frontend UI", "React 18, Vite, Tailwind CSS", "High-fidelity, glassmorphic client interface"],
        ["Animations", "Framer Motion", "Smooth screen transformations, pulse alarms, road lights"],
        ["GIS / Maps", "Leaflet.js, OpenStreetMap", "Cyber GIS mapping, dynamic vector overlays, routes"],
        ["Icons & Assets", "Lucide React, Inline vector SVGs", "High-fidelity diagnostics indicators, active mascots"],
        ["AI Models", "Google Vision API & Gemini", "Pothole scanning, severity indexing, AI dialog bubbles"]
    ]

    for row_idx, data in enumerate(tech_data, start=1):
        for col_idx, text in enumerate(data):
            cell = table.cell(row_idx, col_idx)
            cell.text = text
            set_cell_margins(cell)
            if row_idx % 2 == 0:
                set_cell_background(cell, "F1F5F9")
            else:
                set_cell_background(cell, "FFFFFF")

    doc.add_page_break()

    # ==========================================
    # CHAPTER 8: DEPLOYMENT ARCHITECTURE & FUTURE SCOPE
    # ==========================================
    add_chapter_title("7. Deployment & Future Roadmap")
    add_section_heading("7.1 Deployment Infrastructure")
    doc.add_paragraph(
        "RaastaSense V2 is deployed to a high-availability, edge CDN hosted via Vercel for fast assets delivery. The API telemetry server "
        "is active on Render with custom environmental variables configurations. The live, production application is accessible at:"
    )
    p_link = doc.add_paragraph()
    p_link.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run_l = p_link.add_run("Live Deployed Production Link: https://raastasense.vercel.app")
    run_l.font.bold = True
    run_l.font.size = DocxPt(12)
    run_l.font.color.rgb = C_BLUE

    add_section_heading("7.2 Future Scope")
    doc.add_paragraph("• Edge Computer Vision: Direct integration of real-time custom TensorFlow Lite pothole detectors on citizen smartphones.")
    doc.add_paragraph("• Municipal Blockchain Integration: Smart-contract backed audit trails to automatically award bounty scores and payouts to citizens who submit confirmed hazard points.")
    doc.add_paragraph("• Active V2X Signaling: Integrated cellular vehicle-to-everything (V2X) beacons to alert passing drivers of active SOS crash vectors.")
    doc.add_page_break()

    # ==========================================
    # CHAPTER 9: CONCLUSION
    # ==========================================
    add_chapter_title("8. Conclusion")
    doc.add_paragraph(
        "RaastaSense V2 represents a robust, hackathon-winning paradigm in municipal road safety. By bridging GIS live intelligence, "
        "state-specific legal compliance surcharges, automated AI incident diagnosis, and sat-locked responder beacon overlays, "
        "the NXT-waVE team has built a fail-safe companion ready for real-world municipal deployment. RaastaSense V2 ensures every traveler's "
        "journey is protected by active safety intelligence, demonstrating the powerful impact of smart road safety systems in preserving human life."
    )

    doc.save("RaastaSense_Hackathon_Submission.docx")
    print("DOCX submission report generated successfully!")

if __name__ == "__main__":
    create_presentation()
    create_docx_proposal()
    print("All documents generated beautifully!")
